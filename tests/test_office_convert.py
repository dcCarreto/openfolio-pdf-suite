"""Testes de ConvertOfficeToPDF.

Os testes de fallback forcam find_libreoffice() a retornar None via
monkeypatch, para exercitar sempre o caminho em Python puro
independente de a maquina que roda os testes ter o LibreOffice
instalado. Os testes "via_libreoffice" rodam de verdade quando o
LibreOffice esta disponivel, e sao pulados quando nao esta.
"""

import subprocess

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfReader

import core.office_convert as office_convert
from core.office_convert import ConvertOfficeToPDF, find_libreoffice


def _make_docx(path):
    document = Document()
    document.add_heading("Título do documento", level=1)
    document.add_paragraph("Primeiro parágrafo de teste.")
    document.save(path)


def _make_xlsx(path):
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Dados"
    sheet.append(["Nome", "Valor"])
    sheet.append(["Item 1", 100])
    workbook.save(path)


def _make_pptx(path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Título do slide"
    slide.placeholders[1].text_frame.text = "Conteúdo do slide de teste"
    presentation.save(path)


def _normalized_text(page) -> str:
    # O extract_text() do pypdf pode inserir espaços/quebras espúrias ao redor
    # de caracteres acentuados, mesmo quando o PDF renderiza corretamente
    # (confirmado visualmente durante o desenvolvimento, inclusive com o
    # LibreOffice fazendo a conversão de verdade). Normaliza espaços em
    # branco antes de comparar.
    return " ".join(page.extract_text().split())


def test_docx_to_pdf_fallback(tmp_path, monkeypatch):
    monkeypatch.setattr(office_convert, "find_libreoffice", lambda: None)
    docx_path = tmp_path / "doc.docx"
    _make_docx(docx_path)

    output_path = tmp_path / "doc.pdf"
    ConvertOfficeToPDF().run(str(docx_path), str(output_path))

    reader = PdfReader(str(output_path))
    text = _normalized_text(reader.pages[0])
    assert "Título do documento" in text
    assert "Primeiro" in text and "de teste." in text


def test_xlsx_to_pdf_fallback(tmp_path, monkeypatch):
    monkeypatch.setattr(office_convert, "find_libreoffice", lambda: None)
    xlsx_path = tmp_path / "planilha.xlsx"
    _make_xlsx(xlsx_path)

    output_path = tmp_path / "planilha.pdf"
    ConvertOfficeToPDF().run(str(xlsx_path), str(output_path))

    reader = PdfReader(str(output_path))
    assert "Item 1" in _normalized_text(reader.pages[0])


def test_pptx_to_pdf_fallback(tmp_path, monkeypatch):
    monkeypatch.setattr(office_convert, "find_libreoffice", lambda: None)
    pptx_path = tmp_path / "apresentacao.pptx"
    _make_pptx(pptx_path)

    output_path = tmp_path / "apresentacao.pdf"
    ConvertOfficeToPDF().run(str(pptx_path), str(output_path))

    reader = PdfReader(str(output_path))
    assert len(reader.pages) == 1
    assert "Título do slide" in _normalized_text(reader.pages[0])


def test_rejects_unsupported_format(tmp_path):
    bogus_path = tmp_path / "arquivo.txt"
    bogus_path.write_text("conteudo", encoding="utf-8")

    with pytest.raises(ValueError):
        ConvertOfficeToPDF().run(str(bogus_path), str(tmp_path / "out.pdf"))


def test_libreoffice_conversion_wraps_called_process_error(tmp_path, monkeypatch):
    monkeypatch.setattr(office_convert, "find_libreoffice", lambda: "soffice")
    docx_path = tmp_path / "doc.docx"
    _make_docx(docx_path)

    def fake_run(args, **kwargs):
        raise subprocess.CalledProcessError(1, args, stderr=b"erro simulado do soffice")

    monkeypatch.setattr(office_convert.subprocess, "run", fake_run)

    with pytest.raises(RuntimeError, match="erro simulado do soffice"):
        ConvertOfficeToPDF().run(str(docx_path), str(tmp_path / "out.pdf"))


def test_libreoffice_conversion_wraps_timeout(tmp_path, monkeypatch):
    monkeypatch.setattr(office_convert, "find_libreoffice", lambda: "soffice")
    docx_path = tmp_path / "doc.docx"
    _make_docx(docx_path)

    def fake_run(args, **kwargs):
        raise subprocess.TimeoutExpired(cmd=args, timeout=120)

    monkeypatch.setattr(office_convert.subprocess, "run", fake_run)

    with pytest.raises(RuntimeError, match="timeout"):
        ConvertOfficeToPDF().run(str(docx_path), str(tmp_path / "out.pdf"))


def test_find_libreoffice_does_not_crash():
    result = find_libreoffice()
    assert result is None or isinstance(result, str)


_HAS_LIBREOFFICE = find_libreoffice() is not None
_SKIP_REASON = "LibreOffice não está instalado nesta máquina"


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason=_SKIP_REASON)
def test_docx_to_pdf_via_libreoffice(tmp_path):
    docx_path = tmp_path / "doc.docx"
    _make_docx(docx_path)

    output_path = tmp_path / "doc.pdf"
    ConvertOfficeToPDF().run(str(docx_path), str(output_path))

    reader = PdfReader(str(output_path))
    assert len(reader.pages) >= 1
    assert "Título do documento" in _normalized_text(reader.pages[0])


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason=_SKIP_REASON)
def test_xlsx_to_pdf_via_libreoffice(tmp_path):
    xlsx_path = tmp_path / "planilha.xlsx"
    _make_xlsx(xlsx_path)

    output_path = tmp_path / "planilha.pdf"
    ConvertOfficeToPDF().run(str(xlsx_path), str(output_path))

    reader = PdfReader(str(output_path))
    assert "Item 1" in _normalized_text(reader.pages[0])


@pytest.mark.skipif(not _HAS_LIBREOFFICE, reason=_SKIP_REASON)
def test_pptx_to_pdf_via_libreoffice(tmp_path):
    pptx_path = tmp_path / "apresentacao.pptx"
    _make_pptx(pptx_path)

    output_path = tmp_path / "apresentacao.pdf"
    ConvertOfficeToPDF().run(str(pptx_path), str(output_path))

    reader = PdfReader(str(output_path))
    assert len(reader.pages) >= 1
    assert "Título do slide" in _normalized_text(reader.pages[0])
