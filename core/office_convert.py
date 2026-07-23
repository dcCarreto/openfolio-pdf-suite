"""Conversão de documentos do Office (Word, Excel, PowerPoint) para PDF.

Usa o LibreOffice (se estiver instalado na máquina) para uma conversão fiel
ao arquivo original. Quando o LibreOffice não é encontrado, cai para uma
conversão básica em Python puro (via python-docx/openpyxl/python-pptx +
reportlab), que preserva o texto e a estrutura, mas não a formatação visual
exata do arquivo original.
"""

import shutil
import subprocess
import tempfile
from pathlib import Path
from xml.sax.saxutils import escape

from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.pdfgen import canvas
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

from .base import PDFOperation

_LIBREOFFICE_CANDIDATES = [
    "soffice",
    "libreoffice",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
]

_SUPPORTED_SUFFIXES = (".docx", ".xlsx", ".xls", ".pptx")


def find_libreoffice() -> str | None:
    """Procura o executável do LibreOffice no PATH e em locais comuns de instalação."""
    for candidate in _LIBREOFFICE_CANDIDATES:
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
        if Path(candidate).is_file():
            return candidate
    return None


class ConvertOfficeToPDF(PDFOperation):
    """Converte um documento do Office (.docx, .xlsx, .pptx) para PDF."""

    def run(self, input_path: str, output_path: str) -> None:
        suffix = Path(input_path).suffix.lower()
        if suffix not in _SUPPORTED_SUFFIXES:
            raise ValueError(f"Formato não suportado: {suffix}")

        libreoffice_path = find_libreoffice()
        if libreoffice_path:
            self._convert_with_libreoffice(libreoffice_path, input_path, output_path)
        else:
            self._convert_with_fallback(input_path, output_path, suffix)

    def _convert_with_libreoffice(self, libreoffice_path: str, input_path: str, output_path: str) -> None:
        with tempfile.TemporaryDirectory() as tmp_dir:
            subprocess.run(
                [
                    libreoffice_path,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp_dir,
                    input_path,
                ],
                check=True,
                capture_output=True,
                timeout=120,
            )
            generated_path = Path(tmp_dir) / f"{Path(input_path).stem}.pdf"
            if not generated_path.is_file():
                raise RuntimeError("O LibreOffice não gerou o arquivo PDF esperado.")
            shutil.copy(generated_path, output_path)

    def _convert_with_fallback(self, input_path: str, output_path: str, suffix: str) -> None:
        if suffix == ".docx":
            _docx_to_pdf(input_path, output_path)
        elif suffix in (".xlsx", ".xls"):
            _xlsx_to_pdf(input_path, output_path)
        else:
            _pptx_to_pdf(input_path, output_path)


def _docx_to_pdf(input_path: str, output_path: str) -> None:
    from docx import Document

    document = Document(input_path)
    styles = getSampleStyleSheet()
    story = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            story.append(Spacer(1, 10))
            continue
        style = styles["Heading2"] if paragraph.style.name.startswith("Heading") else styles["Normal"]
        story.append(Paragraph(escape(text), style))
        story.append(Spacer(1, 4))

    if not story:
        story.append(Paragraph("(documento vazio)", styles["Normal"]))

    SimpleDocTemplate(output_path, pagesize=A4).build(story)


def _xlsx_to_pdf(input_path: str, output_path: str) -> None:
    from openpyxl import load_workbook

    workbook = load_workbook(input_path, data_only=True)
    styles = getSampleStyleSheet()
    story = []

    for sheet in workbook.worksheets:
        story.append(Paragraph(escape(sheet.title), styles["Heading2"]))
        rows = [
            ["" if cell is None else str(cell) for cell in row]
            for row in sheet.iter_rows(values_only=True)
        ]
        if rows:
            table = Table(rows)
            table.setStyle(
                TableStyle(
                    [
                        ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                        ("FONTSIZE", (0, 0), (-1, -1), 8),
                        ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ]
                )
            )
            story.append(table)
        story.append(Spacer(1, 20))

    if not story:
        story.append(Paragraph("(planilha vazia)", styles["Normal"]))

    SimpleDocTemplate(output_path, pagesize=landscape(A4)).build(story)


def _pptx_to_pdf(input_path: str, output_path: str) -> None:
    from pptx import Presentation

    presentation = Presentation(input_path)
    width, height = landscape(A4)
    c = canvas.Canvas(output_path, pagesize=(width, height))
    margin = 50

    for slide in presentation.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        texts.append(text)

        y = height - margin
        if texts:
            c.setFont("Helvetica-Bold", 20)
            c.drawString(margin, y, texts[0][:90])
            y -= 36
            c.setFont("Helvetica", 12)
            for text in texts[1:]:
                for line in _wrap_text(text, 95):
                    if y < margin:
                        break
                    c.drawString(margin, y, line)
                    y -= 18
        c.showPage()

    c.save()


def _wrap_text(text: str, width: int) -> list[str]:
    import textwrap

    return textwrap.wrap(text, width=width) or [""]
